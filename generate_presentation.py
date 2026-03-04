#!/usr/bin/env python3
"""
Script to generate PowerPoint presentation for API Monitoring Tool
Requires: pip install python-pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed.")
    print("Please run: pip install python-pptx")
    exit(1)

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    purple = RGBColor(102, 126, 234)  # #667eea
    violet = RGBColor(118, 75, 162)   # #764ba2
    dark_gray = RGBColor(51, 51, 51)
    light_gray = RGBColor(136, 136, 136)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    left = top = Inches(1)
    width = height = Inches(8.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "API Monitoring Tool"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = purple
    p.alignment = PP_ALIGN.CENTER
    
    subtitle = title_frame.add_paragraph()
    subtitle.text = "Dual-Mode Monitoring Solution"
    subtitle.font.size = Pt(32)
    subtitle.font.color.rgb = violet
    subtitle.alignment = PP_ALIGN.CENTER
    
    subtitle2 = title_frame.add_paragraph()
    subtitle2.text = "Azure Application Insights & On-Premise Elasticsearch"
    subtitle2.font.size = Pt(24)
    subtitle2.font.color.rgb = light_gray
    subtitle2.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Overview"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What We Built:"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    
    bullets = [
        "Azure Mode: Monitor APIs hosted in Azure using Application Insights",
        "On-Premise Mode: Monitor APIs hosted locally using Elasticsearch",
        "Unified Dashboard: Single dashboard for both monitoring modes",
        "Centralized Logging: TelemetryAPI handles all data collection"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # Slide 3: Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Architecture Overview"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = """Generic API (On-Premise)
    ↓ HTTP POST /api/telemetry/log/elastic
TelemetryAPI
    ↓                    ↓
Elasticsearch    Application Insights
    ↓                    ↓
    └─────────┬─────────┘
              ↓
         Dashboard"""
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.name = "Courier New"
    
    # Slide 4: Azure Application Insights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Azure Application Insights Setup"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "How It Works:"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    
    steps = [
        "Other APIs (Azure-hosted) send telemetry to TelemetryAPI",
        "TelemetryAPI receives data at /api/telemetry/log",
        "Application Insights stores the data in Azure",
        "Dashboard queries Application Insights using KQL queries"
    ]
    
    for i, step in enumerate(steps, 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {step}"
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Data Flow: Azure API → TelemetryAPI → Application Insights → Dashboard"
    p.level = 0
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = violet
    
    # Slide 5: On-Premise Elasticsearch
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "On-Premise Elasticsearch Setup"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "How It Works:"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    
    steps = [
        "GenericAPI (On-Premise) sends telemetry to TelemetryAPI",
        "TelemetryAPI receives data at /api/telemetry/log/elastic",
        "Elasticsearch stores the data locally",
        "Dashboard queries Elasticsearch directly via REST API"
    ]
    
    for i, step in enumerate(steps, 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {step}"
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Data Flow: GenericAPI → TelemetryAPI → Elasticsearch → Dashboard"
    p.level = 0
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = violet
    
    # Slide 6: TelemetryAPI
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "TelemetryAPI - Centralized Logging"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Purpose:"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    
    bullets = [
        "Single Point for all telemetry data collection",
        "Dual Destination routing based on API type",
        "Separation of Concerns: Logging vs. Querying"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "\nEndpoints:"
    p.font.size = Pt(28)
    p.font.bold = True
    p.space_before = Pt(12)
    
    endpoints = [
        "/api/telemetry/log → Application Insights (Azure APIs)",
        "/api/telemetry/log/elastic → Elasticsearch (GenericAPI)"
    ]
    
    for endpoint in endpoints:
        p = tf.add_paragraph()
        p.text = f"• {endpoint}"
        p.level = 0
        p.font.size = Pt(16)
        p.font.name = "Courier New"
        p.space_after = Pt(8)
    
    # Slide 7: Data Capture - GenericAPI
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Data Capture - GenericAPI (On-Premise)"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Middleware Integration"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• ApiLoggingMiddleware intercepts all HTTP requests"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "• Captures: Method, Path, Status Code, Response Time, Request/Response Body"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "• Extracts: Customer, TenantId from query/headers"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\n2. Telemetry Sending"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "await httpClient.PostAsync(\n  'http://localhost:5292/api/telemetry/log/elastic',\n  jsonContent\n);"
    p.level = 0
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    
    # Slide 8: Data Capture - Azure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Data Capture - Azure APIs"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. API Integration"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• APIs call TelemetryAPI endpoint"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "• Send telemetry data via HTTP POST"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\n2. Application Insights Storage"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Data stored as Custom Events"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "• Event Name: 'ApiCallTelemetry'"
    p.level = 0
    p.font.size = Pt(16)
    p.font.name = "Courier New"
    
    p = tf.add_paragraph()
    p.text = "• Custom Dimensions contain all telemetry data"
    p.level = 0
    p.font.size = Pt(16)
    
    # Slide 9: Dashboard Query Methods
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Dashboard - Query Methods"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Application Insights Mode:"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Query API: Application Insights REST API"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "• Query Language: KQL (Kusto Query Language)"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "• Authentication: API Key"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "\nOn-Premise Mode:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "• Query API: Elasticsearch REST API (Direct)"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "• Query Language: Elasticsearch Query DSL"
    p.level = 0
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "• Authentication: None (local development)"
    p.level = 0
    p.font.size = Pt(16)
    
    # Slide 10: Dashboard Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Dashboard Features"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Metrics:"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    
    metrics = [
        "📊 Total Requests",
        "⏱️ Average Response Time",
        "✅ Success Rate",
        "❌ Error Count"
    ]
    
    for metric in metrics:
        p = tf.add_paragraph()
        p.text = metric
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "\nCharts & Analytics:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_before = Pt(12)
    
    charts = [
        "📈 Request Volume Over Time",
        "🎯 Status Code Breakdown",
        "⏱️ Response Time Distribution",
        "🔝 Top 10 Endpoints"
    ]
    
    for chart in charts:
        p = tf.add_paragraph()
        p.text = chart
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    # Slide 11: Configuration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Configuration Files"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "config.js"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "APP_INSIGHTS_CONFIG: applicationId, apiKey, queryApiUrl"
    p.level = 0
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    
    p = tf.add_paragraph()
    p.text = "ON_PREMISE_CONFIG: elasticsearchUrl, indexName, currentCustomer"
    p.level = 0
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    
    p = tf.add_paragraph()
    p.text = "\nappsettings.json (TelemetryAPI)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "ApplicationInsights: ConnectionString"
    p.level = 0
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    
    p = tf.add_paragraph()
    p.text = "Elasticsearch: Url, IndexName"
    p.level = 0
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    
    # Slide 12: Elasticsearch Setup
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Elasticsearch Setup"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Installation Steps:"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    
    steps = [
        "Download Elasticsearch (8.19.7)",
        "Configure elasticsearch.yml: cluster.name, node.name, network.host, CORS",
        "Start Elasticsearch: START_ELASTICSEARCH.bat",
        "Verify: http://localhost:9200"
    ]
    
    for i, step in enumerate(steps, 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {step}"
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # Slide 13: Complete Data Flow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Complete Data Flow Diagram"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = """GENERIC API (On-Premise IIS)
    ↓ ApiLoggingMiddleware
    ↓ POST /api/telemetry/log/elastic
TELEMETRY API
    ↓                    ↓
Elasticsearch    Application Insights
    ↓                    ↓
    └─────────┬─────────┘
              ↓
         DASHBOARD
    (Browser - Direct Query)"""
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.name = "Courier New"
    
    # Slide 14: Key Differences
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank for table
    left = top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Key Differences"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = purple
    
    # Create table - need rows + 1 for header
    rows = 7  # 1 header + 6 data rows
    cols = 3
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    table.cell(0, 0).text = "Aspect"
    table.cell(0, 1).text = "Application Insights"
    table.cell(0, 2).text = "Elasticsearch"
    
    # Data rows
    data = [
        ("Location", "Azure Cloud", "On-Premise (Local)"),
        ("Storage", "Azure Application Insights", "Local Elasticsearch"),
        ("Query Method", "KQL (Kusto)", "Elasticsearch Query DSL"),
        ("Authentication", "API Key", "None (Local)"),
        ("Endpoint", "/api/telemetry/log", "/api/telemetry/log/elastic"),
        ("Cost", "Pay-per-use", "Free (Self-hosted)")
    ]
    
    for i, (aspect, ai, es) in enumerate(data, 1):
        if i < rows:  # Safety check
            table.cell(i, 0).text = aspect
            table.cell(i, 1).text = ai
            table.cell(i, 2).text = es
    
    # Style header
    for col in range(cols):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = purple
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
    
    # Style data cells
    for row in range(1, rows):
        for col in range(cols):
            cell = table.cell(row, col)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                if col == 0:
                    paragraph.font.bold = True
    
    # Slide 15: Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Benefits"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    benefits = [
        "✅ Centralized Logging - Single TelemetryAPI handles all logging",
        "✅ Flexible Architecture - Support for both Azure and On-Premise",
        "✅ Real-Time Monitoring - Live dashboard updates, auto-refresh",
        "✅ Comprehensive Metrics - Request/Response tracking, Error monitoring",
        "✅ Cost-Effective - On-Premise: Free, Azure: Pay-per-use"
    ]
    
    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = benefit
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # Slide 16: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = purple
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    summary_points = [
        "1. Dual-Mode Monitoring - Azure Application Insights & Elasticsearch",
        "2. Centralized Logging - TelemetryAPI as single point of collection",
        "3. Unified Dashboard - Single interface for both modes",
        "4. Complete Solution - Data capture, storage, and visualization"
    ]
    
    for point in summary_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(16)
    
    # Slide 17: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left = top = Inches(1)
    width = height = Inches(8.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Thank You!"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = purple
    p.alignment = PP_ALIGN.CENTER
    
    subtitle = title_frame.add_paragraph()
    subtitle.text = "Questions?"
    subtitle.font.size = Pt(32)
    subtitle.font.color.rgb = violet
    subtitle.alignment = PP_ALIGN.CENTER
    
    subtitle2 = title_frame.add_paragraph()
    subtitle2.text = "Project: API Monitoring Tool\nTechnologies: ASP.NET Core, Elasticsearch, Application Insights, JavaScript"
    subtitle2.font.size = Pt(18)
    subtitle2.font.color.rgb = light_gray
    subtitle2.alignment = PP_ALIGN.CENTER
    subtitle2.space_before = Pt(40)
    
    return prs

if __name__ == "__main__":
    print("Creating PowerPoint presentation...")
    prs = create_presentation()
    output_file = "API_Monitoring_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation created successfully: {output_file}")
    print(f"Location: {output_file}")

